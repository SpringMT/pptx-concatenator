"""
PowerPoint concatenation library using pptx-slide-copier.
Concatenates multiple PPTX files by appending target slides to source presentation.
"""

from pptx import Presentation
from pptx_slide_copier import SlideCopier
from typing import Union, List
from pathlib import Path


class PptxConcatenator:
    """
    A utility class for concatenating PowerPoint presentations.
    """

    @staticmethod
    def concat(src: Union[str, Path, Presentation],
               target: Union[str, Path, Presentation],
               output_path: Union[str, Path] = None) -> Presentation:
        """
        Concatenate target PPTX to the end of source PPTX.

        Args:
            src: Source presentation (file path or Presentation object)
            target: Target presentation to append (file path or Presentation object)
            output_path: Optional path to save the result. If None, returns Presentation object only.

        Returns:
            Presentation object with concatenated slides

        Example:
            >>> from pptx_concat import PptxConcatenator
            >>> result = PptxConcatenator.concat("source.pptx", "target.pptx", "output.pptx")
        """
        # Load presentations if paths are provided
        if isinstance(src, (str, Path)):
            src_prs = Presentation(str(src))
        else:
            src_prs = src

        if isinstance(target, (str, Path)):
            target_prs = Presentation(str(target))
        else:
            target_prs = target

        # Copy all slides from target to source
        num_target_slides = len(target_prs.slides)
        for slide_idx in range(num_target_slides):
            SlideCopier.copy_slide(target_prs, slide_idx, src_prs)

        # Save if output path is provided
        if output_path:
            src_prs.save(str(output_path))

        return src_prs

    @staticmethod
    def concat_multiple(src: Union[str, Path, Presentation],
                        targets: List[Union[str, Path, Presentation]],
                        output_path: Union[str, Path] = None) -> Presentation:
        """
        Concatenate multiple target PPTX files to the end of source PPTX.

        Args:
            src: Source presentation (file path or Presentation object)
            targets: List of target presentations to append
            output_path: Optional path to save the result

        Returns:
            Presentation object with all concatenated slides

        Example:
            >>> from pptx_concat import PptxConcatenator
            >>> result = PptxConcatenator.concat_multiple(
            ...     "source.pptx",
            ...     ["target1.pptx", "target2.pptx"],
            ...     "output.pptx"
            ... )
        """
        # Load source presentation if path is provided
        if isinstance(src, (str, Path)):
            src_prs = Presentation(str(src))
        else:
            src_prs = src

        # Concatenate each target
        for target in targets:
            if isinstance(target, (str, Path)):
                target_prs = Presentation(str(target))
            else:
                target_prs = target

            # Copy all slides from this target to source
            num_target_slides = len(target_prs.slides)
            for slide_idx in range(num_target_slides):
                SlideCopier.copy_slide(target_prs, slide_idx, src_prs)

        # Save if output path is provided
        if output_path:
            src_prs.save(str(output_path))

        return src_prs


# Convenience function for simple concatenation
def concat_pptx(src: Union[str, Path],
                target: Union[str, Path],
                output_path: Union[str, Path]) -> None:
    """
    Simple function to concatenate two PPTX files.

    Args:
        src: Source PPTX file path
        target: Target PPTX file path to append
        output_path: Output file path

    Example:
        >>> from pptx_concat import concat_pptx
        >>> concat_pptx("source.pptx", "target.pptx", "result.pptx")
    """
    PptxConcatenator.concat(src, target, output_path)
